from django.shortcuts import render
from django.http import HttpResponse
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import csv
from datetime import datetime, timedelta
import matplotlib.pyplot as plt
import io
import os
import tempfile
import matplotlib
matplotlib.use('Agg')
from django import forms

class ReportForm(forms.Form):
    csv_file = forms.FileField(label='Select CSV file', required=True)
    pptx_template = forms.FileField(label='Select PowerPoint template', required=True)
    no_process = forms.BooleanField(label='No Procesar Salidas', required=False)
    nombre_evento = forms.CharField(label='Nombre del Evento', max_length=100, required=True)
    pase_nombre = forms.CharField(label='Nombre del Pase', max_length=100, required=True)
    fecha_evento = forms.DateField(label='Fecha del Evento', required=True)

def power_report(request):
    return render(request, 'generate_report.html')


    
def cargar_archivo_csv(csv_file, pptx_template, no_process, nombre_evento, pase_nombre, fecha_evento, request):
    print("Procesando archivo CSV...")
    csv_file = request.FILES.get('csv_file', None)
    pptx_template = request.FILES.get('pptx_template', None)  # Obtener la plantilla PPTX desde la solicitud
    print("Archivo CSV:", csv_file)
    print("Plantilla PPTX:", pptx_template)
    
    if csv_file:
        archivo_pptx = procesar_archivo_csv(csv_file, pptx_template, no_process, nombre_evento, pase_nombre, fecha_evento, request)  # Pasar los argumentos necesarios
        if isinstance(archivo_pptx, HttpResponse):
            return archivo_pptx
        else:
            response = HttpResponse(archivo_pptx.getvalue(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = 'attachment; filename="informe.pptx"'
            archivo_pptx.close()
            return response
    else:
        return HttpResponse("Error: No se proporcionó ningún archivo CSV.")

    
def procesar_archivo_csv(csv_file, pptx_template , no_process , nombre_evento, pase_nombre, fecha_evento, request):
    print("Procesando archivo CSV...")
    if pptx_template:
        prs = Presentation(pptx_template)  # Usar la plantilla PPTX proporcionada
    else:
        prs = Presentation()  # Si no se proporciona una plantilla, crear una presentación vacía

    csv_content = csv_file.read().decode('latin1').splitlines()
    lines = csv_content[7:]

    if len(lines) < 3:
        return HttpResponse("Error: El archivo CSV no tiene suficientes líneas.")

    no_process = request.POST.get('no_process', None)
    nombre_evento = request.POST.get('nombre_evento', None)
    pase_nombre = request.POST.get('pase_nombre', None)
    fecha_evento = request.POST.get('fecha_evento', None)

    reader = csv.DictReader(lines, delimiter=';')

    puertas_unicas = {}
    categorias_unicas = {}
    suma_tipo_acceso = 0
    fechas_grafico = []  # Inicializar la variable aquí
    entradas = []  # Inicializar la variable aquí
    salidas = []  # Inicializar la variable aquí
    try:
        for row in reader:
            puerta = row['Puerta']
            tipo_acceso_str = row['Tipo acceso']

            if puerta not in puertas_unicas:
                puertas_unicas[puerta] = 0
            puertas_unicas[puerta] += 1

            categoria = row['Categoría']
            if categoria not in categorias_unicas:
                categorias_unicas[categoria] = 0
            categorias_unicas[categoria] += 1

            if tipo_acceso_str.strip():
                if tipo_acceso_str == "Entrada" or tipo_acceso_str == "Salida":
                    suma_tipo_acceso += 1

        if len(prs.slides) >= 3:
            slide3 = prs.slides[2]
        else:
            while len(prs.slides) < 3:
                prs.slides.add_slide(prs.slide_layouts[0])
            slide3 = prs.slides[2]
        # Resto del código para el slide 3

        if len(prs.slides) >= 4:
            slide4 = prs.slides[3]
        else:
            while len(prs.slides) < 4:
                prs.slides.add_slide(prs.slide_layouts[0])
            slide4 = prs.slides[3]
        # Resto del código para el slide 4

        if len(prs.slides) >= 5:
            slide5 = prs.slides[4]
        else:
            while len(prs.slides) < 5:
                prs.slides.add_slide(prs.slide_layouts[0])
            slide5 = prs.slides[4]
        # Resto del código para el slide 5

        slide_width = Inches(10)
        slide_height = Inches(5.625)
        title_height = Inches(0.78)
        table_top = title_height + Inches(0.5)
        available_height = slide_height - title_height - Inches(0.5)

        num_rows = len(puertas_unicas) + 2
        num1_rows = len(categorias_unicas) + 2

        table_width = Inches(8)
        table_height = min(Inches(4), available_height)
        table_left = (slide_width - table_width) / 2

        shape3 = slide3.shapes.add_table(num_rows, 2, table_left, table_top, table_width, table_height).table
        shape4 = slide4.shapes.add_table(num1_rows, 2, table_left, table_top, table_width, table_height).table

        shape3.cell(0, 0).text = "PUERTAS"
        shape3.cell(0, 1).text = "INGRESOS"
        shape4.cell(0, 0).text = "CATEGORIAS"
        shape4.cell(0, 1).text = "INGRESOS"

        for i, puerta in enumerate(puertas_unicas.keys()):
            shape3.cell(i + 1, 0).text = puerta
            shape3.cell(i + 1, 1).text = str(puertas_unicas[puerta])
        for a, categoria in enumerate(categorias_unicas.keys()):
            shape4.cell(a + 1, 0).text = categoria
            shape4.cell(a + 1, 1).text = str(categorias_unicas[categoria])

        total_row3 = shape3.rows[num_rows - 1]
        total_row4 = shape4.rows[num1_rows - 1]

        for cell in total_row3.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(86, 146, 178)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

        for cell in total_row4.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(86, 146, 178)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

        total_row3.cells[0].text = "TOTAL"
        total_row4.cells[0].text = "TOTAL"

        total_ingresos = sum(puertas_unicas.values())
        total_row3.cells[1].text = str(total_ingresos)
        total_ingresos = sum(categorias_unicas.values())
        total_row4.cells[1].text = str(total_ingresos)

        temp_file_path = None
        try:
            temp_file = tempfile.NamedTemporaryFile(delete=False, mode='w', encoding='latin1')
            temp_file_path = temp_file.name
            temp_file.writelines('\n'.join(lines))
            temp_file.close()

            with open(temp_file_path, "r", encoding="latin1") as file:
                for _ in range(0):
                    next(file)
                print(file)
                reader = csv.DictReader(file, delimiter=';')
                fechas = []
                tipo_acceso = []
                for row in reader:
                    fecha_str = row.get('Fecha', None)
                    tipo_acceso_str = row.get('Tipo acceso', None)
                    if fecha_str and tipo_acceso_str and (tipo_acceso_str == "Entrada" or tipo_acceso_str == "Salida"):
                        try:
                            fecha = datetime.strptime(fecha_str, "%d/%m/%Y %H:%M:%S")
                            fechas.append(fecha)
                            tipo_acceso.append(tipo_acceso_str)
                          # Contar las entradas y salidas
                            if tipo_acceso_str == "Entrada":
                               entradas.append(fecha)  # Agregar la fecha de entrada
                            elif tipo_acceso_str == "Salida" and not no_process:
                               salidas.append(fecha)  # Agregar la fecha de salida

                        except ValueError as e:
                            print(f"Error: {e}")
                if fechas:
                    fecha_inicial = min(fechas)
                    fecha_final = max(fechas)
                    fechas_grafico = [fecha_inicial + timedelta(minutes=30 * i) for i in
                                      range(int((fecha_final - fecha_inicial).total_seconds() / 1800) + 1)]
                    entradas = [0] * len(fechas_grafico)
                    salidas = [0] * len(fechas_grafico)
                    for i, fecha in enumerate(fechas_grafico):
                         for j, fecha_csv in enumerate(fechas):
                              if fecha <= fecha_csv < fecha + timedelta(minutes=30):
                                if tipo_acceso[j] == "Entrada":
                                   entradas[i] += 1
                                elif tipo_acceso[j] == "Salida":
                                   salidas[i] += 1

                else:
                    # Mostrar un mensaje de error o registrar la falta de fechas válidas
                    print("No se encontraron fechas válidas en el archivo CSV.")
                print("Datos del gráfico:")
                print("Fechas:", fechas_grafico)
                print("Entradas:", entradas)
                print("Salidas:", salidas)

        except ValueError as e:
            print(f"Error: {e}")
        finally:
            if temp_file_path:
                os.unlink(temp_file_path)

# Resto del código para agregar el gráfico a la presentación
        with io.BytesIO() as image_stream:
            plt.plot(fechas_grafico, entradas, color='orange', label='Entradas')
            plt.plot(fechas_grafico, salidas, color='skyblue', label='Salidas')
            plt.xlabel('Fecha y Hora')
            plt.ylabel('Cantidad')
            plt.title('Ingresos y Salidas por Intervalo de 30 minutos')
            plt.legend()
            plt.xticks(rotation=45)
            plt.tight_layout()
            plt.savefig(image_stream, format='png')
            image_stream.seek(0)
            print("Proceso de visualización:")

            left = Inches(1)
            top = Inches(1)
            slide5.shapes.add_picture(image_stream, left, top, width=Inches(8), height=Inches(4.5))
            print("El gráfico se ha agregado correctamente a la presentación de PowerPoint.")

            plt.close()
        # Resto del código para actualizar los placeholders en la presentación PowerPoint

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text == "{{NOMBRE_EVENTO}}":
                                run.text = nombre_evento
                            elif run.text == "{{PASE_NOMBRE}}":
                                run.text = pase_nombre
                            elif run.text == "{{FECHA_EVENTO}}":
                                run.text = fecha_evento
                            elif run.text == "{{INGRESO_NUMERO}}":
                                run.text = str(suma_tipo_acceso)
        # Guardar la presentación PowerPoint en memoria y retornarla como respuesta HTTP

        pptx_in_memory = io.BytesIO()
        prs.save(pptx_in_memory)
        pptx_in_memory.seek(0)
        return pptx_in_memory

    except Exception as e:
        print(f"Error al procesar el archivo CSV: {e}")
        return HttpResponse("Error al procesar el archivo CSV. Por favor, inténtalo de nuevo.")


def generate_report(request):
    if request.method == 'POST':
        form = ReportForm(request.POST, request.FILES)
        if form.is_valid():
            csv_file = request.FILES['csv_file']
            pptx_template = request.FILES['pptx_template']
            no_process = form.cleaned_data['no_process']
            nombre_evento = form.cleaned_data['nombre_evento']
            pase_nombre = form.cleaned_data['pase_nombre']
            fecha_evento = form.cleaned_data['fecha_evento']
            archivo_pptx = cargar_archivo_csv(csv_file, pptx_template, no_process, nombre_evento, pase_nombre, fecha_evento, request)
            # Resto del código
        if isinstance(archivo_pptx, HttpResponse):
            return archivo_pptx
        else:
            return render(request, 'generate_report.html', {'informe_generado': archivo_pptx})
    else:
        return render(request, 'generate_report.html')
